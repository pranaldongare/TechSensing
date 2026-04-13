import asyncio
import io
import os
from pathlib import Path
import re
import shutil
import time
import traceback
import uuid
import xml.etree.ElementTree as ET

from bs4 import BeautifulSoup
from docx import Document as DocxDocument
from docx.opc.constants import RELATIONSHIP_TYPE as RT
import fitz
import markdown
import olefile
import pandas as pd
from PIL import Image
from pptx import Presentation

from app.socket_handler import sio
from core.config import settings
from core.models.document import Document, Page
from core.parsers.excel_utils import (
    deduplicate_columns,
    detect_merged_header_rows,
    enrich_dataframe_with_metadata,
    find_header_row,
    flatten_multiindex_columns,
)
from core.parsers.extensions import IMAGE_EXTENSIONS, SUPPORTED_EXTENSIONS
from core.parsers.image import image_parser
from core.parsers.slide_export import (
    _timeout_for_file,
    convert_ppt_to_pptx,
    get_libreoffice_command,
)
from core.parsers.vlm import vlm_parse_concurrent, vlm_parse_slide
from core.constants import SWITCHES
from core.parsers.glm_ocr import glm_ocr_parse, glm_ocr_parse_concurrent
from core.services.sqlite_manager import SQLiteManager

try:
    from docling.document_converter import DocumentConverter
    _HAS_DOCLING = True
except ImportError:
    _HAS_DOCLING = False

BASE_DIR = os.path.dirname(os.path.abspath(__file__))
PROJECT_ROOT = os.path.dirname(BASE_DIR)
SMARTART_URI = "http://schemas.openxmlformats.org/drawingml/2006/diagram"

# Regex to extract the label text from a mermaid node definition: A[label text]
_MERMAID_LABEL_RE = re.compile(r'\[([^\]]+)\]')


def _sanitize_mermaid(raw: str) -> str:
    """Remove repetitive VLM-hallucinated node chains from mermaid diagrams."""
    if not raw:
        return raw
    cleaned = raw.strip()
    if cleaned.startswith("```mermaid"):
        cleaned = cleaned[len("```mermaid"):].strip()
    if cleaned.endswith("```"):
        cleaned = cleaned[:-3].strip()

    lines = cleaned.split("\n")
    kept_lines = []
    label_counts: dict = {}
    trimmed = False

    for line in lines:
        labels = _MERMAID_LABEL_RE.findall(line)
        hit_limit = False
        for label in labels:
            normalized = label.strip().lower()
            if not normalized:
                continue
            if label_counts.get(normalized, 0) >= 2:
                hit_limit = True
                break
        if hit_limit:
            trimmed = True
            break
        kept_lines.append(line)
        for label in labels:
            normalized = label.strip().lower()
            if normalized:
                label_counts[normalized] = label_counts.get(normalized, 0) + 1

    result = "\n".join(kept_lines)
    if trimmed:
        print(
            f"[Mermaid] Removed repetitive nodes: kept {len(kept_lines)}/{len(lines)} lines "
            f"({len(result)} chars from {len(cleaned)})"
        )
    return result


XML_NAMESPACES = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "dgm": "http://schemas.openxmlformats.org/drawingml/2006/diagram",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
}


def _clean_ppt_text(text: str) -> str:
    if not text:
        return ""
    return re.sub(r"\s+", " ", text).strip()


def _extract_shapes_recursive(shapes, slide_part, depth=0) -> list[str]:
    """Recursively extract text from all shapes, including GroupShape containers."""
    texts = []
    for shape in shapes:
        try:
            # Recurse into grouped shapes (flowcharts, org charts, etc.)
            if hasattr(shape, "shapes"):
                texts.extend(
                    _extract_shapes_recursive(shape.shapes, slide_part, depth + 1)
                )
                continue

            table_block = _extract_table_block(shape)
            if table_block:
                texts.append(table_block)
                continue

            smartart_block = _extract_smartart_block(shape, slide_part)
            if smartart_block:
                texts.append(smartart_block)
                continue

            if hasattr(shape, "text"):
                cleaned_text = _clean_ppt_text(getattr(shape, "text", ""))
                if cleaned_text:
                    texts.append(cleaned_text)
        except Exception:
            traceback.print_exc()
    return texts


def _extract_table_block(shape) -> str:
    if not getattr(shape, "has_table", False):
        return ""

    try:
        table_rows = []
        for row in shape.table.rows:
            row_cells = [_clean_ppt_text(cell.text) for cell in row.cells]
            if any(row_cells):
                table_rows.append(" | ".join(row_cells))

        if not table_rows:
            return ""

        table_text = "\n".join(table_rows)
        return f"[Table]\n{table_text}\n[/Table]"
    except Exception:
        traceback.print_exc()
        return ""


def _extract_smartart_text_from_xml(xml_blob: bytes) -> list[str]:
    try:
        root = ET.fromstring(xml_blob)
    except Exception:
        traceback.print_exc()
        return []

    texts = []
    for node in root.findall(".//a:t", XML_NAMESPACES):
        cleaned = _clean_ppt_text(node.text or "")
        if cleaned:
            texts.append(cleaned)
    return texts


def _extract_smartart_block(shape, slide_part) -> str:
    try:
        graphic_data = shape.element.find(".//a:graphicData", XML_NAMESPACES)
        if graphic_data is None or graphic_data.get("uri") != SMARTART_URI:
            return ""

        smartart_lines = []

        # Text can exist inline in the shape XML for some decks.
        for node in shape.element.findall(".//a:t", XML_NAMESPACES):
            cleaned = _clean_ppt_text(node.text or "")
            if cleaned:
                smartart_lines.append(cleaned)

        # For native SmartArt, text is often in related diagram parts.
        rel_ids = graphic_data.find(".//dgm:relIds", XML_NAMESPACES)
        if rel_ids is not None:
            rel_keys = (
                "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}dm",
                "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}qs",
                "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}cs",
                "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}lo",
            )
            for rel_key in rel_keys:
                rel_id = rel_ids.get(rel_key)
                if not rel_id:
                    continue

                try:
                    rel = slide_part.rels[rel_id]
                except Exception:
                    continue

                target_part = getattr(rel, "target_part", None)
                xml_blob = getattr(target_part, "blob", None)
                if not xml_blob:
                    continue
                smartart_lines.extend(_extract_smartart_text_from_xml(xml_blob))

        # Deduplicate while preserving order.
        deduped_lines = list(dict.fromkeys(smartart_lines))
        if not deduped_lines:
            return ""

        return (
            "[SmartArt Diagram]\n" + "\n".join(deduped_lines) + "\n[/SmartArt Diagram]"
        )
    except Exception:
        traceback.print_exc()
        return ""


def extract_text_from_doc(path: str) -> str:
    """Extract readable text from a legacy .doc file (pure Python)."""
    if not olefile.isOleFile(path):
        raise ValueError(f"{path} is not a valid .doc file")

    with olefile.OleFileIO(path) as ole:
        if not ole.exists("WordDocument"):
            raise ValueError("No WordDocument stream found")
        stream = ole.openstream("WordDocument")
        data = stream.read()

    # Decode binary to text (best effort)
    text = data.decode("latin-1", errors="ignore")
    # Remove control characters
    text = re.sub(r"[\x00-\x08\x0B\x0C\x0E-\x1F]+", " ", text)
    # Collapse extra whitespace
    text = re.sub(r"\s{2,}", " ", text)
    # Keep only readable ASCII chunks
    text = "\n".join(re.findall(r"[ -~]{5,}", text))
    return text.strip()


async def extract_document(
    path, title="Untitled", file_name=None, user_id=None, thread_id=None
):
    start_time = time.time()
    file_path = path
    ext = Path(path).suffix.lower()
    original_type = ext[1:]  # Preserve original format (e.g. "docx", "pptx") for Document.type
    _converted_pdf_path = None  # Track temp PDF from format conversion for cleanup

    # Derive a safe base name even if file_name is None
    try:
        safe_file_name = file_name or os.path.basename(file_path)
        name, _ = os.path.splitext(safe_file_name)
    except Exception:
        traceback.print_exc()
        safe_file_name = os.path.basename(file_path)
        name, _ = os.path.splitext(safe_file_name)

    # Normalize user/thread to avoid crashing on None
    user_id = user_id or "unknown_user"
    thread_id = thread_id or "unknown_thread"
    doc_id = str(uuid.uuid4())

    async def safe_emit(channel: str, payload: dict):
        try:
            await sio.emit(channel, payload)
        except Exception as e:
            print(f"[emit-error] channel={channel} payload={payload} err={e}")

    if ext not in SUPPORTED_EXTENSIONS:
        print(f"Unsupported file type: {ext} for {safe_file_name}. Skipping.")
        await safe_emit(
            f"{user_id}/progress",
            {"message": f"Skipping {title}: unsupported file type {ext}"},
        )
        return None

    # --- Handle standalone images ---
    if ext in IMAGE_EXTENSIONS:
        try:
            await safe_emit(
                f"{user_id}/progress",
                {"message": f"{title} is an image, extracting text..."},
            )
            text = await image_parser(file_path)
        except Exception as e:
            print(f"Error processing image {safe_file_name}: {str(e)}")
            traceback.print_exc()
            return None

        # --- GLM-OCR Enhancement for standalone images (additive) ---
        if SWITCHES.get("GLM_OCR", False):
            try:
                print(f"[Image] Running GLM-OCR enhancement on {safe_file_name}...")
                glm_result = await glm_ocr_parse(file_path, mode="text")
                if glm_result and glm_result.strip():
                    text += f"\n\n{glm_result.strip()}"
                    print(
                        f"[Image] GLM-OCR enhancement added ({len(glm_result)} chars)"
                    )
            except Exception as e:
                print(f"[Image] GLM-OCR enhancement failed: {e}")
                traceback.print_exc()

        await safe_emit(
            f"{user_id}/progress",
            {"message": f"Processed {safe_file_name} successfully"},
        )

        end_time = time.time()
        return Document(
            id=doc_id,
            type=ext[1:],
            file_name=safe_file_name,
            content=[Page(number=1, text=text)],
            title=title,
            full_text=text,
        )

    # --- Handle Markdown files ---
    if ext == ".md":
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                md_text = f.read()

            # Convert markdown -> HTML -> plain text
            try:
                html = markdown.markdown(md_text)
            except Exception:
                traceback.print_exc()
                html = md_text  # fallback
            try:
                soup = BeautifulSoup(html, "html.parser")
                plain_text = soup.get_text(separator="\n")
            except Exception:
                traceback.print_exc()
                plain_text = md_text

            # Prepare image handling (keyed by doc_id to avoid filename collisions)
            image_dir = f"data/{user_id}/threads/{thread_id}/images/{doc_id}"
            try:
                os.makedirs(image_dir, exist_ok=True)
            except Exception:
                traceback.print_exc()

            ocr_tasks = {}
            image_names = []

            # Regex to find Markdown image syntax: ![alt](path)
            image_pattern = re.compile(r"!\[.*?\]\((.*?)\)")
            matches = image_pattern.findall(md_text)

            page_text = plain_text
            for idx, img_path in enumerate(matches, start=1):
                try:
                    resolved_path = img_path
                    if not os.path.isabs(resolved_path):
                        # make path relative to md file
                        resolved_path = os.path.join(
                            os.path.dirname(file_path), resolved_path
                        )

                    if not os.path.exists(resolved_path):
                        print(f"Markdown image not found: {resolved_path}")
                        continue

                    ext_img = Path(resolved_path).suffix.lstrip(".")
                    image_name = f"md_img{idx}.{ext_img}"
                    dest_path = os.path.join(image_dir, image_name)

                    # Copy image into project folder
                    try:
                        shutil.copy(resolved_path, dest_path)
                    except Exception:
                        traceback.print_exc()
                        continue
                    image_names.append(image_name)

                    placeholder = f"{{PENDING_{image_name}}}"
                    page_text += f"\n\n{placeholder}"

                    # Run OCR asynchronously
                    ocr_tasks[placeholder] = asyncio.create_task(
                        image_parser(dest_path)
                    )
                except Exception:
                    traceback.print_exc()

            # Wait for OCR tasks
            for placeholder, task in ocr_tasks.items():
                try:
                    image_text = await task
                except Exception as e:
                    print(f"Error parsing Markdown image: {e}")
                    traceback.print_exc()
                    image_text = "[Image OCR failed]"

                page_text = page_text.replace(placeholder, image_text, 1)

            await safe_emit(
                f"{user_id}/progress",
                {"message": f"Processed {safe_file_name} (Markdown) successfully"},
            )

            return Document(
                id=doc_id,
                type="markdown",
                file_name=safe_file_name,
                content=[Page(number=1, text=page_text, images=image_names)],
                title=title,
                full_text=md_text,  # preserve original markdown
            )

        except Exception as e:
            print(f"Error processing Markdown file {safe_file_name}: {str(e)}")
            traceback.print_exc()
            return None

    if ext in {".xls", ".xlsx", ".csv"}:
        try:
            # Read Excel or CSV file into DataFrame(s)
            # Read Excel or CSV file into DataFrame(s)
            sheets_data = {}  # Tuples of (df, context_text)

            if ext == ".xlsx":
                # Use robust parsing for modern Excel
                xls = pd.ExcelFile(file_path, engine="openpyxl")
                for sheet_name in xls.sheet_names:
                    # 1. Detect Header & Context
                    header_idx, context = find_header_row(file_path, sheet_name)

                    # 2. Detect multi-level headers from merged cells
                    header_param = detect_merged_header_rows(
                        file_path, sheet_name, header_idx
                    )

                    # 3. Read DataFrame with correct header(s)
                    df = pd.read_excel(xls, sheet_name=sheet_name, header=header_param)

                    # 4. Flatten MultiIndex columns if multi-level headers detected
                    if isinstance(header_param, list):
                        df = flatten_multiindex_columns(df)

                    # 5. Enrich with Metadata (Colors, Comments)
                    # Note: We pass header_idx so we know where data starts
                    enrichment_header = (
                        header_param[-1]
                        if isinstance(header_param, list)
                        else header_param
                    )
                    df = enrich_dataframe_with_metadata(
                        df, file_path, sheet_name, enrichment_header
                    )

                    sheets_data[sheet_name] = (df, context)

            elif ext == ".xls":
                # Legacy Excel (less features supported, no openpyxl enrichment)
                xls = pd.ExcelFile(file_path, engine="xlrd")
                for sheet_name in xls.sheet_names:
                    # Heuristic for header might still work if we adapt it for xlrd,
                    # but current utility uses openpyxl.
                    # Fallback to standard read for .xls to avoid complexity deps
                    df = pd.read_excel(xls, sheet_name=sheet_name)
                    sheets_data[sheet_name] = (df, None)

            else:
                # CSV
                df = pd.read_csv(file_path)
                sheets_data["Sheet1"] = (df, None)

            # --- Load into SQLite for structured querying ---
            try:
                tables_info = SQLiteManager.load_spreadsheet(
                    user_id=user_id,
                    thread_id=thread_id,
                    doc_id=doc_id,
                    file_path=file_path,
                    file_name=safe_file_name,
                )
                if tables_info:
                    print(
                        f"[SQLite] Loaded {len(tables_info)} table(s) for {safe_file_name}: "
                        f"{list(tables_info.keys())}"
                    )
            except Exception as e:
                print(f"[SQLite] Failed to load {safe_file_name} into SQLite: {e}")
                traceback.print_exc()

            # --- Generate Global Workbook Context (Phase 3) ---
            # Create a summary of all sheets to give the LLM a "Table of Contents"
            workbook_summary_lines = ["# Workbook Structure Summary"]
            for s_name, (s_df, _) in sheets_data.items():
                col_list = ", ".join(
                    [str(c) for c in s_df.columns[:10]]
                )  # Limit cols to avoid huge headers
                if len(s_df.columns) > 10:
                    col_list += ", ..."
                workbook_summary_lines.append(
                    f"- Sheet '{s_name}': {len(s_df)} rows. Columns: [{col_list}]"
                )

            workbook_summary = "\n".join(workbook_summary_lines) + "\n\n"

            # --- Also generate text representation for RAG/vector store ---
            text_parts = []
            pages = []
            page_num = 1

            for sheet_name, (df, context_text) in sheets_data.items():
                # Drop fully empty rows
                df = df.dropna(how="all")

                # --- Clean unicode whitespace (non-breaking spaces etc.) ---
                # Apply to ALL columns: replace \u00a0 and other unicode whitespace
                for col in df.columns:
                    if df[col].dtype == object or str(df[col].dtype) == "string":
                        df[col] = df[col].apply(
                            lambda x: (
                                re.sub(
                                    r"[\u00a0\u200b\u200c\u200d\ufeff\xa0]+",
                                    " ",
                                    str(x),
                                )
                                .replace("\n", " ")
                                .strip()
                                if isinstance(x, str) and str(x) != "nan"
                                else x
                            )
                        )

                # --- Fix "Unnamed" columns ---
                # Replace 'Unnamed: N' column headers with something more useful
                new_cols = []
                for i, col in enumerate(df.columns):
                    col_str = str(col)
                    # Clean non-breaking spaces from column names too
                    col_str = re.sub(r"[\u00a0\xa0]+", " ", col_str).strip()
                    if col_str.startswith("Unnamed"):
                        # Try to use the first non-null value in the column as a hint
                        first_val = df.iloc[:, i].dropna().head(1)
                        if not first_val.empty:
                            hint = str(first_val.iloc[0]).strip()
                            hint = re.sub(r"[\u00a0\xa0]+", " ", hint).strip()
                            # Only use as column name if it looks like a header (short text)
                            if (
                                hint
                                and len(hint) < 50
                                and not hint.replace(".", "").replace(",", "").isdigit()
                            ):
                                col_str = hint
                            else:
                                col_str = f"Column_{i}"
                        else:
                            col_str = f"Column_{i}"
                    new_cols.append(col_str)
                df.columns = new_cols

                # Deduplicate column names (handles duplicates after cleanup)
                df.columns = deduplicate_columns(list(df.columns))

                # Drop columns that are entirely NaN or empty
                df = df.dropna(axis=1, how="all")

                # Replace NaN with empty string for cleaner text output
                df = df.fillna("")

                # Remove rows where all values are empty strings
                df = df[
                    df.apply(lambda row: any(str(v).strip() != "" for v in row), axis=1)
                ]

                # Build a text summary: schema + data
                col_info = ", ".join([str(col) for col in df.columns])

                # Prepend the context (pre-header text) if it exists
                context_block = ""
                if context_text:
                    context_block = f"Context/Metadata:\n{context_text}\n"

                sheet_header = (
                    f"=== Spreadsheet: {safe_file_name} | Sheet: {sheet_name} ===\n"
                    f"{workbook_summary}"  # <--- Global Context
                    f"{context_block}"  # <--- Local Sheet Context
                    f"Columns: {col_info}\n"
                    f"Total rows: {len(df)}\n"
                )

                # Use markdown table for RAG — much more readable for the LLM
                try:
                    data_text = df.to_markdown(index=False)
                except Exception:
                    # Fallback: try to_string which is still more readable than JSON
                    try:
                        data_text = df.to_string(index=False)
                    except Exception:
                        data_text = str(df)

                # Final cleanup: remove any remaining \u00a0 from the output
                data_text = re.sub(r"[\u00a0\xa0]+", " ", data_text)

                sheet_text = sheet_header + "\nData:\n" + data_text
                # Collapse excessive whitespace but preserve single newlines for table rows
                sheet_text = re.sub(r"[^\S\n]{2,}", " ", sheet_text).strip()

                text_parts.append(sheet_text)
                pages.append(Page(number=page_num, text=sheet_text))
                page_num += 1

            full_text = "\n\n".join(text_parts)

            # Get the schema info to store with the document
            schema = SQLiteManager.get_schema(user_id, thread_id)

            await safe_emit(
                f"{user_id}/progress",
                {"message": f"Processed {safe_file_name} (Excel/CSV) successfully"},
            )

            return Document(
                id=doc_id,
                type="spreadsheet",
                file_name=safe_file_name,
                content=pages,
                title=title,
                full_text=full_text,
                has_sql_data=True,
                spreadsheet_schema=schema,
            )

        except Exception as e:
            print(f"Error processing Excel/CSV file {safe_file_name}: {str(e)}")
            traceback.print_exc()
            return None

    # --- Handle legacy Word .doc files — convert to PDF, fall through to PDF handler ---
    if ext == ".doc":
        await safe_emit(
            f"{user_id}/progress",
            {"message": f"Converting {title} (.doc) to modern format..."},
        )

        libreoffice_cmd = get_libreoffice_command()

        # Try DOC → PDF directly (preferred — feeds into unified PDF pipeline)
        if libreoffice_cmd:
            try:
                doc_dir = os.path.dirname(file_path)
                proc = await asyncio.create_subprocess_exec(
                    libreoffice_cmd,
                    "--headless",
                    "--convert-to",
                    "pdf",
                    "--outdir",
                    doc_dir,
                    file_path,
                    stdout=asyncio.subprocess.PIPE,
                    stderr=asyncio.subprocess.PIPE,
                )
                await asyncio.wait_for(proc.communicate(), timeout=120)
                if proc.returncode == 0:
                    expected_pdf = os.path.splitext(file_path)[0] + ".pdf"
                    if os.path.exists(expected_pdf):
                        _converted_pdf_path = expected_pdf
                        file_path = expected_pdf
                        ext = ".pdf"
                        print(f"[DOC] Converted to PDF, delegating to PDF pipeline")
            except Exception as e:
                print(f"[DOC] LibreOffice DOC→PDF conversion failed: {e}")
                traceback.print_exc()

        # Fallback: DOC → DOCX (which will then try PDF conversion in DOCX handler)
        if ext == ".doc" and libreoffice_cmd:
            try:
                doc_dir = os.path.dirname(file_path)
                proc = await asyncio.create_subprocess_exec(
                    libreoffice_cmd,
                    "--headless",
                    "--convert-to",
                    "docx",
                    "--outdir",
                    doc_dir,
                    file_path,
                    stdout=asyncio.subprocess.PIPE,
                    stderr=asyncio.subprocess.PIPE,
                )
                await asyncio.wait_for(proc.communicate(), timeout=120)
                if proc.returncode == 0:
                    expected = os.path.splitext(file_path)[0] + ".docx"
                    if os.path.exists(expected):
                        try:
                            os.remove(file_path)
                        except Exception:
                            pass
                        file_path = expected
                        ext = ".docx"
                        print(f"[DOC] Converted to DOCX, delegating to DOCX handler")
            except Exception as e:
                print(f"[DOC] LibreOffice DOC→DOCX conversion failed: {e}")
                traceback.print_exc()
        # Last resort: binary extraction if all conversions failed
        if ext == ".doc":
            try:
                text = extract_text_from_doc(file_path)
            except Exception as e:
                print(f"Error processing .doc file {safe_file_name}: {str(e)}")
                traceback.print_exc()
                return None

            return Document(
                id=doc_id,
                type="doc",
                file_name=safe_file_name,
                content=[Page(number=1, text=text)],
                title=title,
                full_text=text,
            )

    # --- Handle PowerPoint files — convert to PDF, fall through to PDF handler ---
    if ext in {".ppt", ".pptx"}:

        # If .ppt, convert to .pptx first (needed for python-pptx fallback)
        if ext == ".ppt":
            await safe_emit(
                f"{user_id}/progress",
                {"message": f"Converting {safe_file_name} to modern PPTX format..."},
            )

            converted_path = await convert_ppt_to_pptx(file_path)

            if not converted_path:
                print(f"[Parser] Failed to convert {safe_file_name} to .pptx")
                return None

            try:
                os.remove(file_path)
            except Exception:
                pass

            file_path = converted_path
            ext = ".pptx"

        # Convert PPTX → PDF for unified processing
        libreoffice_cmd = get_libreoffice_command()
        print(f"[PPT] LibreOffice command: {libreoffice_cmd}")
        if libreoffice_cmd:
            try:
                await safe_emit(
                    f"{user_id}/progress",
                    {"message": f"Converting {safe_file_name} to PDF for processing..."},
                )
                ppt_dir = os.path.dirname(file_path)
                lo_timeout = _timeout_for_file(file_path)
                print(f"[PPT] Running LibreOffice: {libreoffice_cmd} --headless --convert-to pdf --outdir {ppt_dir} {file_path} (timeout={lo_timeout}s)")
                proc = await asyncio.create_subprocess_exec(
                    libreoffice_cmd, "--headless", "--convert-to", "pdf",
                    "--outdir", ppt_dir, file_path,
                    stdout=asyncio.subprocess.PIPE, stderr=asyncio.subprocess.PIPE,
                )
                stdout, stderr = await asyncio.wait_for(proc.communicate(), timeout=lo_timeout)
                print(f"[PPT] LibreOffice returncode: {proc.returncode}")
                if stdout and stdout.strip():
                    print(f"[PPT] LibreOffice stdout: {stdout.decode(errors='replace')[:500]}")
                if stderr and stderr.strip():
                    print(f"[PPT] LibreOffice stderr: {stderr.decode(errors='replace')[:500]}")
                # LibreOffice may crash (e.g. SIGABRT=134 from GStreamer/multimedia)
                # AFTER writing the PDF.  Check for the output file regardless of
                # return code — the conversion often succeeds before the crash.
                expected_pdf = os.path.splitext(file_path)[0] + ".pdf"
                if proc.returncode != 0:
                    print(f"[PPT] LibreOffice exited with code {proc.returncode}, checking if PDF was written anyway...")
                if os.path.exists(expected_pdf) and os.path.getsize(expected_pdf) > 0:
                    _converted_pdf_path = expected_pdf
                    file_path = expected_pdf
                    ext = ".pdf"
                    if proc.returncode == 0:
                        print(f"[PPT] Converted to PDF, delegating to PDF pipeline")
                    else:
                        print(f"[PPT] LibreOffice crashed (code {proc.returncode}) but PDF was written — using it")
                else:
                    print(f"[PPT] PDF not found at {expected_pdf}, falling back to python-pptx")
            except asyncio.TimeoutError:
                print(f"[PPT] LibreOffice PDF conversion timed out after {lo_timeout}s for {safe_file_name}")
            except Exception as e:
                print(f"[PPT] LibreOffice PDF conversion failed: {e}")
                traceback.print_exc()

    # Fallback: python-pptx extraction if PPTX → PDF conversion failed
    if ext in {".ppt", ".pptx"}:
        try:
            prs = Presentation(file_path)
        except Exception as e:
            print(f"Error opening presentation {safe_file_name}: {e}")
            traceback.print_exc()
            return None

        pages = []
        combined_texts = []
        ocr_tasks = {}
        image_dir = f"data/{user_id}/threads/{thread_id}/images/{doc_id}"
        try:
            os.makedirs(image_dir, exist_ok=True)
        except Exception:
            traceback.print_exc()

        try:
            for slide_number, slide in enumerate(prs.slides, start=1):
                # Extract text recursively (handles GroupShape for flowcharts etc.)
                slide_text = _extract_shapes_recursive(slide.shapes, slide.part)
                page_text = "\n".join(slide_text)

                image_names = []

                # Extract images
                for shape_index, shape in enumerate(slide.shapes, start=1):
                    try:
                        if getattr(shape, "shape_type", None) == 13:  # PICTURE
                            image = shape.image
                            image_bytes = image.blob
                            image_ext = image.ext
                            image_name = (
                                f"slide{slide_number}_img{shape_index}.{image_ext}"
                            )
                            image_path = os.path.join(image_dir, image_name)

                            try:
                                with open(image_path, "wb") as f:
                                    f.write(image_bytes)
                            except Exception:
                                traceback.print_exc()
                                continue

                            placeholder = f"{{PENDING_{image_name}}}"
                            page_text += f"\n\n{placeholder}"
                            image_names.append(image_name)

                            ocr_tasks[placeholder] = asyncio.create_task(
                                image_parser(image_path)
                            )
                    except Exception:
                        traceback.print_exc()

                combined_texts.append(page_text)
                pages.append(
                    Page(number=slide_number, text=page_text, images=image_names)
                )

            # Wait for OCR tasks
            for placeholder, task in ocr_tasks.items():
                try:
                    image_text = await task
                except Exception as e:
                    print(f"Error parsing PPT image: {e}")
                    traceback.print_exc()
                    image_text = "[Image OCR failed]"

                for page in pages:
                    if placeholder in page.text:
                        page.text = page.text.replace(placeholder, image_text, 1)
                combined_texts = [
                    txt.replace(placeholder, image_text, 1) for txt in combined_texts
                ]
        except Exception:
            traceback.print_exc()

        end_time = time.time()
        print(f"Time taken to process {safe_file_name} (.pptx fallback): {end_time - start_time} seconds")
        return Document(
            id=doc_id,
            type=original_type,
            file_name=safe_file_name,
            content=pages,
            title=title,
            full_text="\n".join(combined_texts),
        )

    # --- Handle Word .docx files — convert to PDF, fall through to PDF handler ---
    if ext == ".docx":
        try:
            await safe_emit(
                f"{user_id}/progress",
                {"message": f"Parsing {title} (Word document)..."},
            )

            # Convert DOCX → PDF for unified processing
            libreoffice_cmd = get_libreoffice_command()
            if libreoffice_cmd:
                try:
                    docx_dir = os.path.dirname(file_path)
                    proc = await asyncio.create_subprocess_exec(
                        libreoffice_cmd,
                        "--headless",
                        "--convert-to",
                        "pdf",
                        "--outdir",
                        docx_dir,
                        file_path,
                        stdout=asyncio.subprocess.PIPE,
                        stderr=asyncio.subprocess.PIPE,
                    )
                    await asyncio.wait_for(proc.communicate(), timeout=120)
                    expected_pdf = os.path.splitext(file_path)[0] + ".pdf"
                    if os.path.exists(expected_pdf) and os.path.getsize(expected_pdf) > 0:
                        _converted_pdf_path = expected_pdf
                        file_path = expected_pdf
                        ext = ".pdf"
                        if proc.returncode != 0:
                            print(f"[DOCX] LibreOffice crashed (code {proc.returncode}) but PDF was written — using it")
                        else:
                            print(f"[DOCX] Converted to PDF, delegating to PDF pipeline")
                    elif proc.returncode != 0:
                        print(f"[DOCX] LibreOffice failed (code {proc.returncode}), no PDF produced")
                except Exception as e:
                    print(f"[DOCX] LibreOffice PDF conversion failed: {e}")
                    traceback.print_exc()
        except Exception as e:
            print(f"Error during DOCX conversion: {str(e)}")
            traceback.print_exc()
            return None

    # Fallback: python-docx extraction if DOCX → PDF conversion failed
    if ext == ".docx":
        try:
            print(f"[DOCX] LibreOffice unavailable, using python-docx fallback for {safe_file_name}")
            docx_doc = DocxDocument(file_path)

            image_names_all = []
            ocr_tasks = {}
            image_dir = f"data/{user_id}/threads/{thread_id}/images/{doc_id}"
            try:
                os.makedirs(image_dir, exist_ok=True)
            except Exception:
                traceback.print_exc()

            # --- Extract body elements in document order (paragraphs + tables) ---
            body_parts = []
            for element in docx_doc.element.body:
                tag = element.tag.split("}")[-1] if "}" in element.tag else element.tag
                if tag == "p":
                    # Find the matching Paragraph object
                    for para in docx_doc.paragraphs:
                        if para._element is element:
                            text = para.text.strip()
                            if not text:
                                break
                            style_name = para.style.name if para.style else ""
                            if style_name.startswith("Heading"):
                                try:
                                    level = int(
                                        style_name.replace("Heading ", "").replace(
                                            "Heading", "1"
                                        )
                                    )
                                except (ValueError, TypeError):
                                    level = 1
                                body_parts.append(f"{'#' * level} {text}")
                            else:
                                body_parts.append(text)
                            break
                elif tag == "tbl":
                    # Find the matching Table object
                    for table in docx_doc.tables:
                        if table._element is element:
                            try:
                                table_rows = []
                                # Build header
                                if table.rows:
                                    header_cells = [
                                        cell.text.strip().replace("|", "\\|")
                                        for cell in table.rows[0].cells
                                    ]
                                    table_rows.append(
                                        "| " + " | ".join(header_cells) + " |"
                                    )
                                    table_rows.append(
                                        "| "
                                        + " | ".join(["---"] * len(header_cells))
                                        + " |"
                                    )
                                    # Data rows
                                    for row in table.rows[1:]:
                                        row_cells = [
                                            cell.text.strip().replace("|", "\\|")
                                            for cell in row.cells
                                        ]
                                        table_rows.append(
                                            "| " + " | ".join(row_cells) + " |"
                                        )
                                if table_rows:
                                    body_parts.append(
                                        f"[Table]\n"
                                        + "\n".join(table_rows)
                                        + f"\n[/Table]"
                                    )
                            except Exception:
                                traceback.print_exc()
                            break

            page_text = "\n\n".join(body_parts)

            # --- Extract embedded images ---
            img_index = 0
            MIN_IMAGE_SIZE = 50
            for rel in docx_doc.part.rels.values():
                try:
                    if "image" in rel.reltype:
                        img_index += 1
                        image_part = rel.target_part
                        image_bytes = image_part.blob
                        image_ext = image_part.content_type.split("/")[-1]
                        if image_ext == "jpeg":
                            image_ext = "jpg"

                        img = Image.open(io.BytesIO(image_bytes))
                        if img.width < MIN_IMAGE_SIZE or img.height < MIN_IMAGE_SIZE:
                            continue

                        image_name = f"docx_img{img_index}.{image_ext}"
                        image_path = os.path.join(image_dir, image_name)
                        try:
                            img.save(image_path)
                        except Exception:
                            traceback.print_exc()
                            continue

                        image_names_all.append(image_name)
                        placeholder = f"{{PENDING_{image_name}}}"
                        page_text += f"\n\n{placeholder}"
                        ocr_tasks[placeholder] = asyncio.create_task(
                            image_parser(image_path)
                        )
                except Exception:
                    traceback.print_exc()

            # --- Wait for OCR tasks ---
            for placeholder, task in ocr_tasks.items():
                try:
                    image_text = await task
                except Exception as e:
                    print(f"Error parsing DOCX image: {e}")
                    traceback.print_exc()
                    image_text = "[Image OCR failed]"
                page_text = page_text.replace(placeholder, image_text, 1)

            # Split into pages
            if len(page_text) <= 5000:
                pages = [Page(number=1, text=page_text, images=image_names_all)]
            else:
                sections = page_text.split("\n\n")
                current_page = ""
                page_num = 1
                for section in sections:
                    if len(current_page) + len(section) > 3000 and current_page:
                        pages.append(Page(number=page_num, text=current_page.strip()))
                        page_num += 1
                        current_page = section
                    else:
                        current_page += "\n\n" + section if current_page else section
                if current_page.strip():
                    pages.append(Page(number=page_num, text=current_page.strip(), images=image_names_all))

            end_time = time.time()
            print(f"Time taken to process {safe_file_name} (.docx fallback): {end_time - start_time} seconds")
            return Document(
                id=doc_id,
                type="docx",
                file_name=safe_file_name,
                content=pages,
                title=title,
                full_text=page_text,
            )

        except Exception as e:
            print(f"Error processing DOCX file {safe_file_name}: {str(e)}")
            traceback.print_exc()
            return None

    # --- Handle PDFs and other fitz-supported formats ---
    if ext in [
        ".pdf",
        ".xlsx",
        ".epub",
        ".odt",
        ".txt",
        ".rtf",
        ".html",
        ".xml",
    ]:
        try:
            doc = fitz.open(file_path)
        except Exception as e:
            print(f"Error opening PDF {safe_file_name}: {e}")
            traceback.print_exc()
            return None
        pages = []
        combined_texts = []
        ocr_tasks = {}
        vlm_candidates = []  # Collect pages for concurrent VLM processing
        glm_ocr_candidates = []  # Collect pages for concurrent GLM-OCR processing

        image_dir_base = f"data/{user_id}/threads/{thread_id}/images/{doc_id}"
        MIN_IMAGE_SIZE = 50  # Skip images smaller than 50px (icons, bullets)
        _ocr_xref_cache = {}  # Cache OCR results by image xref to skip duplicates

        # --- Docling native PDF extraction (per-page) ---
        # Docling's export_to_markdown(page_no=N) returns markdown for a single
        # page (1-indexed).  We store the converted document object and call it
        # per page inside the loop so each page gets its own text.
        _docling_doc = None
        if ext == ".pdf" and _HAS_DOCLING:
            try:
                print(f"[PDF] Running Docling extraction for {safe_file_name}...")
                await safe_emit(f"{user_id}/progress", {"message": f"Running Docling extraction for {safe_file_name}..."})

                def _run_docling(p):
                    converter = DocumentConverter()
                    result = converter.convert(p)
                    return result.document

                _docling_doc = await asyncio.to_thread(_run_docling, file_path)
                print(f"[PDF] Docling conversion done — {len(_docling_doc.pages)} pages detected")
            except Exception as e:
                print(f"[PDF] Docling extraction failed: {e}")
                traceback.print_exc()

        for page_number in range(len(doc)):
            page = doc.load_page(page_number)

            table_blocks = []
            page_text = ""
            if _docling_doc:
                try:
                    # Docling uses 1-indexed page numbers
                    docling_page_md = _docling_doc.export_to_markdown(page_no=page_number + 1)
                    if docling_page_md and docling_page_md.strip():
                        page_text = docling_page_md.strip()
                        print(f"[PDF] Docling page {page_number + 1}: {len(page_text)} chars")
                except Exception as e:
                    print(f"[PDF] Docling per-page export failed for page {page_number + 1}: {e}")
                    page_text = ""

            if not page_text:
                # --- Original Table-aware text extraction ---
                try:
                    tables = page.find_tables()
                    for table in tables.tables:
                        try:
                            table_md = table.to_markdown()
                            if table_md and table_md.strip():
                                table_blocks.append(f"[Table]\n{table_md}\n[/Table]")
                        except Exception:
                            pass
                except Exception:
                    pass

                try:
                    table_rects = []
                    try:
                        for table in tables.tables:
                            table_rects.append(fitz.Rect(table.bbox))
                    except Exception:
                        pass

                    if table_rects:
                        text_dict = page.get_text("dict")
                        non_table_lines = []
                        for block in text_dict.get("blocks", []):
                            if block.get("type") != 0:
                                continue
                            block_rect = fitz.Rect(block["bbox"])
                            if not any(block_rect.intersects(tr) for tr in table_rects):
                                for line in block.get("lines", []):
                                    line_text = " ".join(span["text"] for span in line.get("spans", []))
                                    if line_text.strip():
                                        non_table_lines.append(line_text.strip())
                        page_text = "\n".join(non_table_lines)
                    else:
                        page_text = page.get_text("text")

                except Exception:
                    traceback.print_exc()
                    page_text = page.get_text("text")

            # Append table blocks after the regular text
            if table_blocks:
                page_text += "\n\n" + "\n\n".join(table_blocks)

            # --- VLM candidate detection (Phase 1: collect, don't process yet) ---
            # VLM runs on every page when USE_VISION_MODEL=True (the default).
            # Set USE_VISION_MODEL=False in .env to disable (e.g. CPU-only environments).
            if settings.USE_VISION_MODEL:
                try:
                    def _render_pdf(p): return p.get_pixmap(dpi=150).tobytes("png")
                    img_bytes = await asyncio.to_thread(_render_pdf, page)
                    vlm_candidates.append(
                        {
                            "page_index": page_number,
                            "page_number": page_number + 1,
                            "image_bytes": img_bytes,
                            "original_text": page_text,
                            "has_tables": bool(table_blocks),
                        }
                    )
                except Exception as e:
                    print(f"[PDF] Failed to render page {page_number + 1} for VLM: {e}")
                    traceback.print_exc()

            # --- GLM-OCR candidate detection (runs alongside existing OCR) ---
            if SWITCHES.get("GLM_OCR", False):
                try:
                    def _render_pdf_glm(p): return p.get_pixmap(dpi=150).tobytes("png")
                    img_bytes = await asyncio.to_thread(_render_pdf_glm, page)
                    glm_ocr_candidates.append(
                        {
                            "page_index": page_number,
                            "page_number": page_number + 1,
                            "image_bytes": img_bytes,
                        }
                    )
                except Exception as e:
                    print(f"[PDF] Failed to render page {page_number + 1} for GLM-OCR: {e}")
                    traceback.print_exc()

            # Append table blocks after the regular text
            if table_blocks:
                page_text += "\n\n" + "\n\n".join(table_blocks)

            image_names = []
            image_dir = image_dir_base
            try:
                os.makedirs(image_dir, exist_ok=True)
            except Exception:
                traceback.print_exc()

            # Extract embedded raster images and schedule OCR
            try:
                image_list = page.get_images(full=True)
            except Exception:
                traceback.print_exc()
                image_list = []

            for img_index, img in enumerate(image_list):
                try:
                    xref = img[0]
                    base_image = doc.extract_image(xref)
                    image_bytes = base_image.get("image")
                    image_ext = base_image.get("ext", "png")
                    if not image_bytes:
                        continue
                    image = Image.open(io.BytesIO(image_bytes))

                    # Skip tiny decorative images (icons, bullets, logos)
                    if image.width < MIN_IMAGE_SIZE or image.height < MIN_IMAGE_SIZE:
                        continue

                    # Duplicate image detection: reuse OCR results for same xref
                    # (e.g. repeated headers/footers/logos across pages)
                    if xref in _ocr_xref_cache:
                        cached_text = _ocr_xref_cache[xref]
                        if cached_text:
                            page_text += f"\n\n{cached_text}"
                        continue

                    image_name = f"page{page_number + 1}_img{img_index + 1}.{image_ext}"
                    image_path = os.path.join(image_dir, image_name)
                    try:
                        image.save(image_path)
                    except Exception:
                        traceback.print_exc()
                        continue

                    # Put placeholder where the image OCR result should go
                    placeholder = f"{{PENDING_{image_name}}}"
                    page_text += f"\n\n{placeholder}"
                    image_names.append(image_name)

                    # Detect large images that may be diagrams/flowcharts/ERDs
                    # Heuristic: image covers >25% of page area and is at least 400px
                    # in both dimensions — likely a diagram, not a photo or icon.
                    page_area = page.rect.width * page.rect.height
                    img_area = image.width * image.height
                    is_diagram_candidate = (
                        settings.USE_VISION_MODEL
                        and image.width >= 400
                        and image.height >= 400
                        and page_area > 0
                        and (img_area / page_area) > 0.25
                    )

                    if is_diagram_candidate:
                        # Use VLM mermaid prompt for potential diagrams
                        from core.parsers.vlm import vlm_parse_slide

                        async def _mermaid_and_cache(path, xref_id):
                            result = await vlm_parse_slide(path, prompt_type="mermaid")
                            if result:
                                sanitized = _sanitize_mermaid(result)
                                mermaid_text = f"```mermaid\n{sanitized}\n```"
                            else:
                                # Fallback to regular OCR if mermaid extraction fails
                                mermaid_text = await image_parser(path)
                            _ocr_xref_cache[xref_id] = mermaid_text
                            return mermaid_text

                        print(f"[PDF] Page {page_number + 1}: large image {image.width}x{image.height} → Mermaid VLM")
                        ocr_tasks[placeholder] = asyncio.create_task(
                            _mermaid_and_cache(image_path, xref)
                        )
                    else:
                        # OCR raster images, cache result by xref to avoid re-processing duplicates
                        async def _ocr_and_cache(path, xref_id):
                            result = await image_parser(path)
                            _ocr_xref_cache[xref_id] = result
                            return result

                        ocr_tasks[placeholder] = asyncio.create_task(
                            _ocr_and_cache(image_path, xref)
                        )
                except Exception:
                    traceback.print_exc()

            # Extract Diagram/Flowchart bounding boxes (vector graphics) for Mermaid parsing
            try:
                drawings = page.get_drawings()
                # Group connected/overlapping drawings into diagram bounding boxes
                initial_rects = []
                for d in drawings:
                    rect = d["rect"]
                    # Ignore tiny drawing commands (e.g. single lines or dots)
                    if rect.width < 50 or rect.height < 50:
                        continue
                    initial_rects.append(fitz.Rect(rect))

                if initial_rects:
                    print(f"[PDF] Page {page_number + 1}: {len(drawings)} drawings, {len(initial_rects)} pass size filter")

                # Merge overlapping rects (with a margin) to form complete diagram boxes
                diagram_rects = []
                for rect in initial_rects:
                    merged = False
                    for existing in diagram_rects:
                        expanded = fitz.Rect(rect.x0 - 20, rect.y0 - 20, rect.x1 + 20, rect.y1 + 20)
                        if existing.intersects(expanded):
                            existing.include_rect(rect)
                            merged = True
                            break
                    if not merged:
                        diagram_rects.append(rect)

                # Second pass to catch transitive overlaps
                final_rects = []
                for r in diagram_rects:
                    merged = False
                    for existing in final_rects:
                        expanded = fitz.Rect(r.x0 - 20, r.y0 - 20, r.x1 + 20, r.y1 + 20)
                        if existing.intersects(expanded):
                            existing.include_rect(r)
                            merged = True
                            break
                    if not merged:
                        final_rects.append(r)

                # Crop each diagram and schedule for Mermaid VLM extraction
                diagram_count = 0
                for d_index, d_rect in enumerate(final_rects):
                    # Ensure rect is within page bounds
                    d_rect.intersect(page.rect)
                    if d_rect.is_empty or d_rect.width < 100 or d_rect.height < 100:
                        continue

                    diagram_count += 1
                    diagram_name = f"page{page_number + 1}_diagram{d_index + 1}.png"
                    diagram_path = os.path.join(image_dir, diagram_name)

                    # Render just the cropped area
                    try:
                        diagram_pix = page.get_pixmap(clip=d_rect, dpi=150)
                        diagram_bytes = diagram_pix.tobytes("png")

                        image = Image.open(io.BytesIO(diagram_bytes))
                        image.save(diagram_path)
                    except Exception:
                        traceback.print_exc()
                        continue

                    placeholder = f"{{PENDING_{diagram_name}}}"
                    # Stitch surrounding caption context into the diagram block so
                    # caption text ("Figure 2: ...") and diagram content always land
                    # in the same chunk, enabling retrieval by figure number.
                    _ctx = page_text.strip()[-300:] if page_text.strip() else ""
                    _ctx_block = f"[Caption context: ...{_ctx}]\n" if _ctx else ""
                    page_text += f"\n\n[Diagram Detected]\n{_ctx_block}{placeholder}\n[/Diagram Detected]"
                    image_names.append(diagram_name)

                    from core.parsers.vlm import vlm_parse_slide

                    async def _extract_mermaid(path):
                        result = await vlm_parse_slide(path, prompt_type="mermaid")
                        if result:
                            sanitized = _sanitize_mermaid(result)
                            return f"```mermaid\n{sanitized}\n```"
                        return "[Diagram Mermaid Extraction Failed]"

                    ocr_tasks[placeholder] = asyncio.create_task(
                        _extract_mermaid(diagram_path)
                    )
                if diagram_count:
                    print(f"[PDF] Page {page_number + 1}: {diagram_count} vector diagrams queued for Mermaid extraction")
            except Exception:
                traceback.print_exc()

            combined_texts.append(page_text)
            pages.append(
                Page(number=page_number + 1, text=page_text, images=image_names)
            )

        # --- Phase 2: Concurrent VLM processing for candidate pages ---
        if vlm_candidates:
            print(
                f"[PDF] {len(vlm_candidates)} pages queued for concurrent VLM processing"
            )
            await safe_emit(
                f"{user_id}/progress",
                {
                    "message": f"Running VLM on {len(vlm_candidates)} pages of {safe_file_name}..."
                },
            )

            # Split into table-heavy and regular pages for prompt specialization
            table_pages = [c for c in vlm_candidates if c.get("has_tables")]
            regular_pages = [c for c in vlm_candidates if not c.get("has_tables")]

            all_vlm_results = {}  # page_index -> vlm_text

            # Run table pages with table-specific prompt
            if table_pages:
                print(f"[PDF] {len(table_pages)} pages with tables → table prompt")
                table_results = await vlm_parse_concurrent(
                    images=[c["image_bytes"] for c in table_pages],
                    page_labels=[f"Page {c['page_number']}" for c in table_pages],
                    max_concurrent=3,
                    prompt_type="table",
                )
                for c, result in zip(table_pages, table_results):
                    all_vlm_results[c["page_index"]] = result

            # Run regular pages with default prompt
            if regular_pages:
                regular_results = await vlm_parse_concurrent(
                    images=[c["image_bytes"] for c in regular_pages],
                    page_labels=[f"Page {c['page_number']}" for c in regular_pages],
                    max_concurrent=3,
                )
                for c, result in zip(regular_pages, regular_results):
                    all_vlm_results[c["page_index"]] = result

            # Merge VLM results back into pages (always additive)
            # PyMuPDF text layer + VLM visual layer + Mermaid diagrams + GLM-OCR = maximum coverage.
            # Never replace: content-rich pages would lose good PyMuPDF text extraction.
            for candidate in vlm_candidates:
                vlm_text = all_vlm_results.get(candidate["page_index"], "")
                if not vlm_text:
                    continue

                page_idx = candidate["page_index"]
                pages[page_idx].text += f"\n\n{vlm_text}"
                combined_texts[page_idx] += f"\n\n{vlm_text}"
                prompt_used = "table" if candidate.get("has_tables") else "default"
                print(
                    f"[PDF] VLM ({prompt_used}) appended to page {candidate['page_number']} ({len(vlm_text)} chars)"
                )

        # --- Phase 3: Concurrent GLM-OCR processing for candidate pages ---
        if glm_ocr_candidates:
            print(
                f"[PDF] {len(glm_ocr_candidates)} pages queued for GLM-OCR processing"
            )
            await safe_emit(
                f"{user_id}/progress",
                {
                    "message": f"Running GLM-OCR on {len(glm_ocr_candidates)} pages of {safe_file_name}..."
                },
            )

            glm_images = [c["image_bytes"] for c in glm_ocr_candidates]
            glm_labels = [f"Page {c['page_number']}" for c in glm_ocr_candidates]

            glm_results = await glm_ocr_parse_concurrent(
                images=glm_images,
                page_labels=glm_labels,
                mode="text",
                max_concurrent=3,
            )

            # Append GLM-OCR results to existing page content (additive)
            for candidate, glm_text in zip(glm_ocr_candidates, glm_results):
                if not glm_text:
                    continue

                page_idx = candidate["page_index"]
                pages[page_idx].text += f"\n\n{glm_text}"
                combined_texts[page_idx] += f"\n\n{glm_text}"
                print(
                    f"[PDF] GLM-OCR enhancement added to page {candidate['page_number']} ({len(glm_text)} chars)"
                )

        # Wait for OCR tasks from the embedded raster images only
        for placeholder, task in ocr_tasks.items():
            try:
                image_text = await task
            except Exception as e:
                print(f"Error parsing image: {e}")
                traceback.print_exc()
                image_text = "[Image OCR failed]"

            # Replace placeholder once per occurrence (should be exactly 1)
            for page in pages:
                if placeholder in page.text:
                    page.text = page.text.replace(placeholder, image_text, 1)
            combined_texts = [
                txt.replace(placeholder, image_text, 1) for txt in combined_texts
            ]

        # Persist converted PDF for query-time VLM page rendering + cleanup
        if _converted_pdf_path:
            if SWITCHES.get("USE_VLM_FOR_ANSWER", False):
                try:
                    import shutil
                    persist_dir = os.path.join("data", user_id, "threads", thread_id, "parsed")
                    os.makedirs(persist_dir, exist_ok=True)
                    persist_path = os.path.join(persist_dir, f"{doc_id}_pages.pdf")
                    shutil.copy2(_converted_pdf_path, persist_path)
                    print(f"[PDF] Persisted converted PDF → {persist_path}")
                except Exception as e:
                    print(f"[PDF] Failed to persist converted PDF: {e}")
            try:
                os.remove(_converted_pdf_path)
            except Exception:
                pass

        await safe_emit(
            f"{user_id}/progress", {"message": f"Processing {title} successfully..."}
        )
        end_time = time.time()
        print(
            f"Time taken to process {title} successfully: {end_time - start_time} seconds"
        )
        return Document(
            id=doc_id,
            type=original_type,
            file_name=safe_file_name,
            content=pages,
            title=title,
            full_text="\n".join(combined_texts),
        )

    # If we reach here, the extension is supported but not yet implemented
    print(
        f"Parsing for file type {ext} not implemented for {safe_file_name}. Skipping."
    )
    await safe_emit(
        f"{user_id}/progress",
        {"message": f"Skipping {title}: parser for {ext} not implemented yet"},
    )
    return None
